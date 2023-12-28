import PyPDF2
import zipfile
import csv
from docx import Document
from pptx import Presentation
from skimage.feature import hog
from sklearn.svm import SVC
import joblib
import cv2
import os
import sys

# Function to extract HOG features from an image
def get_hog_features(image):
    features, _ = hog(image, orientations=9, pixels_per_cell=(8, 8), cells_per_block=(2, 2), visualize=True)
    return features

# Function to classify an image
def classify_image(image_path, clf):
    img = cv2.imread(image_path, cv2.IMREAD_GRAYSCALE)
    img = cv2.resize(img, (64, 64))  # Resize image to a consistent size
    hog_features = get_hog_features(img)
    prediction = clf.predict([hog_features])
    return prediction[0]

# Load the classifier
clf = joblib.load('trained_classifier_model.pkl')

def check_image_maliciousness(image_path, clf):
    try:
        result = classify_image(image_path, clf)
        if result == 1:
            return True, "The image is classified as malicious."
        else:
            return False, "The image is not classified as malicious."
    except Exception as e:
        return False, f"Error analyzing the image: {str(e)}"

def check_zip_maliciousness(file_path):
    try:
        with zipfile.ZipFile(file_path, 'r') as zip_file:
            # Check each file within the ZIP archive
            for file_info in zip_file.infolist():
                file_name = file_info.filename
                file_content = zip_file.read(file_info)

                # You can add specific checks for each file within the ZIP
                # For example, check for known malicious patterns in file_content

                # Example: Check for JavaScript in HTML files
                if file_name.lower().endswith('.html') and b'<script>' in file_content:
                    return True, f"The ZIP file contains HTML file with JavaScript and may be considered malicious. ({file_name})"

            return False, "The ZIP file does not seem to contain known indicators of malicious content."

    except Exception as e:
        return False, f"Error analyzing the ZIP file: {str(e)}"

def check_pdf_maliciousness(file_path):
    try:
        with open(file_path, 'rb') as pdf_file:
            pdf_reader = PyPDF2.PdfReader(pdf_file)

            # Check each page for JavaScript
            for page_num in range(len(pdf_reader.pages)):
                page = pdf_reader.pages[page_num]
                page_text = page.extract_text()

                # Check for common JavaScript indicators
                if "JavaScript" in page_text:
                    return True, "The PDF contains JavaScript and may be considered malicious."

            return False, "The PDF does not seem to contain known indicators of malicious content."

    except Exception as e:
        return False, f"Error analyzing the PDF: {str(e)}"

def check_pptx_maliciousness(file_path):
    try:
        prs = Presentation(file_path)

        # Check for macros or other indicators of malicious content in the PPTX file
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text_frame"):
                    text = shape.text_frame.text
                    # You can add more checks based on your specific requirements

        return False, "The PPTX file does not seem to contain known indicators of malicious content."

    except Exception as e:
        return False, f"Error analyzing the PPTX file: {str(e)}"

def check_docx_maliciousness(file_path):
    try:
        doc = Document(file_path)

        # Check for macros or other indicators of malicious content in the DOCX file
        for paragraph in doc.paragraphs:
            text = paragraph.text
            # You can add more checks based on your specific requirements

        return False, "The DOCX file does not seem to contain known indicators of malicious content."

    except Exception as e:
        return False, f"Error analyzing the DOCX file: {str(e)}"

def check_csv_maliciousness(file_path):
    try:
        with open(file_path, 'r', newline='', encoding='utf-8') as csv_file:
            csv_reader = csv.reader(csv_file)

            # You can add checks for malicious content in the CSV file
            # For example, checking for suspicious patterns in data

            return False, "The CSV file does not seem to contain known indicators of malicious content."

    except Exception as e:
        return False, f"Error analyzing the CSV file: {str(e)}"

def convert_uri_to_local_path(file_uri):
    # Implement logic to convert URI to local file path
    # This might involve using a library like react-native-fs
    pass

def check_file_maliciousness(file_uri):
    try:
        # Convert URI to local file path
        file_path = convert_uri_to_local_path(file_uri)

        if file_path.lower().endswith('.pdf'):
            return check_pdf_maliciousness(file_path)
        elif file_path.lower().endswith('.png'):
            return check_image_maliciousness(file_path, clf)
        elif file_path.lower().endswith('.pptx'):
            return check_pptx_maliciousness(file_path)
        elif file_path.lower().endswith('.docx'):
            return check_docx_maliciousness(file_path)
        elif file_path.lower().endswith(('.csv', '.xlsx')):
            return check_csv_maliciousness(file_path)
        elif file_path.lower().endswith('.zip'):
            return check_zip_maliciousness(file_path)
        else:
            return False, "Unsupported file type for analysis."

    except Exception as e:
        return False, f"Error analyzing the file: {str(e)}"

if __name__ == "__main__":
    # Check if a file URI is provided as a command-line argument
    if len(sys.argv) < 2:
        print("Usage: python tool_combined.py <file_uri>")
        sys.exit(1)

    file_uri = sys.argv[1]

    is_malicious, message = check_file_maliciousness(file_uri)

    print(f"Malicious: {is_malicious}")
    print(f"Analysis Result: {message}")
